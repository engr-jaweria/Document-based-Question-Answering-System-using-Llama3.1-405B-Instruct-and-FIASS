import os
import tempfile
import numpy as np
import faiss
import requests
import pandas as pd
from sentence_transformers import SentenceTransformer
from langchain.prompts import PromptTemplate
from langchain.chains import RetrievalQA
from langchain_community.vectorstores import FAISS
from langchain.embeddings import HuggingFaceEmbeddings
import streamlit as st
import fitz  # PyMuPDF for PDF handling
from docx import Document as DocxDocument
from pptx import Presentation
from dotenv import load_dotenv
from langchain.schema import Document
from langchain.docstore import InMemoryDocstore
from langchain.vectorstores import FAISS as LangchainFAISS
from langchain.docstore.document import Document as LangchainDocument

# Load environment variables
load_dotenv()

# Set environment variables
REPLICATE_API_TOKEN = os.getenv("REPLICATE_API_TOKEN")
HUGGINGFACE_API_TOKEN = os.getenv("HUGGINGFACE_API_TOKEN")
if not REPLICATE_API_TOKEN or not HUGGINGFACE_API_TOKEN:
    raise ValueError("API tokens for Replicate and Hugging Face are not set in environment variables.")

# Initialize models
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
embeddings_model = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

# Define utility functions
def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file."""
    text = ""
    with fitz.open(pdf_file) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(docx_file):
    """Extract text from a DOCX file."""
    doc = DocxDocument(docx_file)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    """Extract text from a PPTX file."""
    presentation = Presentation(pptx_file)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(xlsx_file):
    """Extract text from an XLSX file."""
    df = pd.read_excel(xlsx_file)
    return df.to_string()

def llama3_1_generate(prompt, model="meta-llama/Llama-3.1-405B-Instruct", top_p=0.9, temperature=0.7, max_tokens=800):
    """Generate text using Llama 3.1."""
    input = {
        "top_p": top_p,
        "prompt": prompt,
        "temperature": temperature,
        "max_tokens": max_tokens,
    }
    # Replace this mock with actual API call or inference logic for Llama 3.1
    return f"Generated response: {prompt}"

# Main Streamlit application
def main():
    st.title("Llama3.1 Document Q&A System")

    # Sidebar inputs
    st.sidebar.header("Upload Input")
    uploaded_files = st.sidebar.file_uploader("Upload documents", type=["txt", "pdf", "docx", "pptx", "xlsx"], accept_multiple_files=True)

    documents = []

    # Validate and process uploaded files
    if uploaded_files:
        temp_dir = tempfile.TemporaryDirectory()

        for uploaded_file in uploaded_files:
            file_path = os.path.join(temp_dir.name, uploaded_file.name)
            with open(file_path, "wb") as f:
                f.write(uploaded_file.read())

            if uploaded_file.name.endswith(".txt"):
                with open(file_path, "r", encoding="utf-8") as f:
                    documents.append(f.read())
            elif uploaded_file.name.endswith(".pdf"):
                documents.append(extract_text_from_pdf(file_path))
            elif uploaded_file.name.endswith(".docx"):
                documents.append(extract_text_from_docx(file_path))
            elif uploaded_file.name.endswith(".pptx"):
                documents.append(extract_text_from_pptx(file_path))
            elif uploaded_file.name.endswith(".xlsx"):
                documents.append(extract_text_from_xlsx(file_path))

        temp_dir.cleanup()

    if documents:
        # Split documents into fixed-size chunks of 1000 characters
        chunk_size = 1000
        chunks = [doc[i:i+chunk_size] for doc in documents for i in range(0, len(doc), chunk_size)]
        docs = [Document(page_content=chunk, metadata={"source": "uploaded file"}) for chunk in chunks]

        # Generate embeddings
        embeddings = embeddings_model.embed_documents(chunks)

        # Convert embeddings to NumPy array (shape: num_embeddings x embedding_dim)
        embeddings_array = np.array(embeddings)
        dimension = embeddings_array.shape[1]
        faiss_index = faiss.IndexFlatL2(dimension)
        faiss_index.add(embeddings_array.astype(np.float32))

        # Create a docstore
        docstore = InMemoryDocstore({i: docs[i] for i in range(len(docs))})

        # Prepare LangChain documents
        vector_store = LangchainFAISS.from_documents(docs, embeddings_model)

        # Set up retriever
        retriever = vector_store.as_retriever(search_kwargs={"k": 5})

        # Chat interface
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []

        st.header("Ask a Question")
        question = st.text_input("Your Question:")

        # Chat interface with better context handling
        if st.button("Get Answer") and question:
            # Integrate chat history into the query
            history_context = " ".join(
                [f"Q: {entry['question']} A: {entry['answer']}" for entry in st.session_state.chat_history]
            )
            prompt = f"Context: {context}\n\nChat History: {history_context}\n\nQuestion: {question}"
            
            # Format prompt based on user input
            if format_choice == "Bullet Points":
                prompt += "\nPlease provide the answer as bullet points."
            elif format_choice == "Summary":
                prompt += "\nPlease summarize concisely."
            elif format_choice == "Specific Length":
                prompt += f"\nLimit the answer to {word_limit} words."
        
            # Generate response using Llama 3.1
            answer = llama3_1_generate(prompt)
        
            # Save chat history and display answer
            st.session_state.chat_history.append({"question": question, "answer": answer})
            st.markdown(f"### **Q: {question}**")
            st.markdown(f"#### **A:** {answer}")
        
            # Show full chat history
            st.header("Chat History")
            for entry in st.session_state.chat_history:
                st.write(f"**Q:** {entry['question']}\n**A:** {entry['answer']}\n")


    else:
        st.write("Upload documents to begin.")

if __name__ == "__main__":
    main()
